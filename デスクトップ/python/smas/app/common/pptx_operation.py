from pptx import Presentation
from pptx.util import Pt, Inches,Cm
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XyChartData

#テンプレート読み込み
def read_template():
	import_dir = "import/pptx/"
	file_name = "template.pptx"
	template_path = import_dir + file_name
	prs = Presentation(template_path)
	return prs

#save
def save_file(prs):
	export_path = "export/"
	file_name = "test.pptx"
	save_path = export_path + file_name
	prs.save(save_path)

#スライド追加
def add_slide(prs):
	#スライドレイアウトの選択
	title_slide_layout = prs.slide_layouts[1]
	#プレゼンテーションにスライドを追加
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	#プレースホルダーの設定
	subtitle = slide.placeholders[1]
	#タイトル、サブタイトルの書き込み
	title.text = "Hello, World!"
	subtitle.text = "python-pptx was here!"
	return slide

#スライドオブジェクトの取得
def get_slide_ob(prs,index):
	slide = prs.slides[index]
	return slide

def get_width(prs):
	width = prs.slide_width
	return width

def get_height(prs):
	height = prs.slide_height
	return height

#画像追加
def add_img(slide):
	export_path = "export/"
	file_name = "気象データ.png"
	png_path = export_path + file_name
	left = slide.shapes[1].left
	top = slide.shapes[1].top
	height = Inches(4.9)
	pic = slide.shapes.add_picture(png_path, left,top,Cm(30), Cm(12))
	return pic

#テーブル追加
def add_tbl(slide,row_nums,col_nums,row_names,col_names,values):
	left = slide.shapes[1].left
	top = slide.shapes[1].top
	table_shape = slide.shapes.add_table(row_nums, col_nums,left,top, Cm(30), Cm(12))
	table = table_shape.table
	#見出し行
	for i in range(len(row_names)):
		cell = table.cell(i+1, 0)
		cell.text = row_names[i]	
	#見出し列
	for i in range(len(col_names)):
		cell = table.cell(0, i+1)
		cell.text = col_names[i]
	#値
	for i in range(len(values)):
		for k in range(len(values[i])):
			cell = table.cell(k+1, i+1)
			cell.text = str(values[i][k])
	return table

#実行テスト
def exec_pptx():
	prs = read_template()
	slide_1 = prs.slides[1]
	slide_2 = prs.slides[2]
	pic = add_img(slide_1)
	row_nums = 6
	col_nums = 6
	row_names = ['カテゴリー１', 'カテゴリー2',  'カテゴリー3']
	col_names = ["a","b","c"]
	values = [[1, 2, 3],[4, 5, 6],[7, 8, 9]]
	table = add_tbl(slide_2,row_nums,col_nums,row_names,col_names,values)
	save_file(prs)

#作成中(グラフ作成)
def create_chart(df):
	x = "平均気温"
	y = "最高気温"
	chart = XyChartData()
	c1 = chart.add_series('系列1')
	c2 = chart.add_series('系列2')
	str_df = df[[x,y]].astype("str")
	df_col = list([str_df.columns.tolist()])
	df_list = df_col + str_df.values.tolist()
	s1.add_data_point()
	print(df_list)

